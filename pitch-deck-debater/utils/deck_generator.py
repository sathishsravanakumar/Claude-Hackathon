"""
PowerPoint deck generator
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from typing import Dict, Any, List
from anthropic import Anthropic
import io
import os
import json


class DeckGenerator:
    """Generate improved PowerPoint presentations"""

    def __init__(self):
        self.client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

    def generate(self, debate_results: Dict[str, Any]) -> bytes:
        """
        Generate an improved PowerPoint deck based on debate results

        Args:
            debate_results: Results from the debate process

        Returns:
            PowerPoint file as bytes
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(prs, debate_results)

        # Executive summary of feedback
        self._add_executive_summary(prs, debate_results)

        # Slide-by-slide improvements
        self._add_improvements_slides(prs, debate_results)

        # Recommendations summary
        self._add_recommendations_slide(prs, debate_results)

        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()

    def _add_title_slide(self, prs: Presentation, debate_results: Dict[str, Any]):
        """Add title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Improved Pitch Deck"
        subtitle.text = "Based on Multi-Agent AI Analysis & Debate"

    def _add_executive_summary(self, prs: Presentation, debate_results: Dict[str, Any]):
        """Add executive summary slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Executive Summary"

        # Add text box with synthesis summary
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        synthesis = debate_results.get("synthesis", "No synthesis available")
        # Extract first 500 characters for summary
        summary_text = synthesis[:500] + "..." if len(synthesis) > 500 else synthesis

        p = text_frame.paragraphs[0]
        p.text = summary_text
        p.font.size = Pt(14)

    def _add_improvements_slides(self, prs: Presentation, debate_results: Dict[str, Any]):
        """Add slides showing improvements for each original slide"""
        original_slides = debate_results.get("deck_content", {}).get("slides", [])
        synthesis = debate_results.get("synthesis", "")

        # Process ALL slides (not just first 5)
        for slide_data in original_slides:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            title = slide.shapes.title
            title.text = f"Slide {slide_data['slide_number']}: {slide_data['title']}"

            # Add improvement notes
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Extract actual improvements from debate for this specific slide
            slide_improvements = self._extract_slide_improvements(
                slide_data,
                synthesis,
                slide_data['slide_number']
            )

            improvements = f"Original Content:\n{slide_data.get('title', 'No title')}\n\n"
            improvements += "Recommended Improvements:\n"
            improvements += slide_improvements

            p = text_frame.paragraphs[0]
            p.text = improvements
            p.font.size = Pt(12)

    def _add_recommendations_slide(self, prs: Presentation, debate_results: Dict[str, Any]):
        """Add final recommendations slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Key Recommendations"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Extract actual recommendations from debate
        recommendations = self._extract_top_recommendations(debate_results)

        p = text_frame.paragraphs[0]
        p.text = recommendations
        p.font.size = Pt(14)

    def _extract_slide_improvements(self, slide_data: Dict[str, Any], synthesis: str, slide_num: int) -> str:
        """Extract specific improvements for a slide from debate synthesis"""
        slide_content = f"Title: {slide_data.get('title', '')}\n"
        slide_content += f"Content: {slide_data.get('content', [])}"

        prompt = f"""Based on the debate synthesis, provide 3-5 specific improvements for Slide {slide_num}.

Slide Content:
{slide_content}

Debate Synthesis:
{synthesis[:1500]}

Return ONLY the bulleted improvements (one per line starting with '-'), no other text."""

        try:
            response = self.client.messages.create(
                model="claude-3-5-haiku-20241022",
                max_tokens=512,
                messages=[{"role": "user", "content": prompt}]
            )
            return response.content[0].text
        except:
            return "- Review and strengthen content\n- Ensure clarity and focus\n- Add supporting data"

    def _extract_top_recommendations(self, debate_results: Dict[str, Any]) -> str:
        """Extract top priority recommendations from debate"""
        synthesis = debate_results.get("synthesis", "")

        prompt = f"""Extract the top 5-7 priority recommendations from this debate synthesis.

Synthesis:
{synthesis}

Return ONLY a numbered list (1., 2., 3., etc.) of actionable recommendations, no other text."""

        try:
            response = self.client.messages.create(
                model="claude-3-5-haiku-20241022",
                max_tokens=1024,
                messages=[{"role": "user", "content": prompt}]
            )
            return response.content[0].text
        except:
            return "1. Strengthen value proposition\n2. Add market validation\n3. Improve financial projections"
