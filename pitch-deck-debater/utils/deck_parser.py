"""
PowerPoint deck parsing utilities
"""
from pptx import Presentation
from typing import List, Dict
import io

def parse_deck(pptx_file) -> List[Dict]:
    """Extract content from PowerPoint deck"""
    if hasattr(pptx_file, 'read'):
        prs = Presentation(io.BytesIO(pptx_file.read()))
    else:
        prs = Presentation(pptx_file)
    
    slides = []
    
    for idx, slide in enumerate(prs.slides):
        slide_data = {
            "number": idx + 1,
            "title": _get_slide_title(slide),
            "content": _extract_slide_text(slide),
            "notes": _get_speaker_notes(slide),
            "shape_count": len(slide.shapes),
            "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown"
        }
        slides.append(slide_data)
    
    return slides

def _get_slide_title(slide) -> str:
    """Extract slide title"""
    if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
        if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
            title_text = slide.shapes.title.text.strip()
            if title_text:
                return title_text
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            return shape.text.strip()[:100]
    
    return "Untitled Slide"

def _extract_slide_text(slide) -> str:
    """Extract all text from slide"""
    text_parts = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_parts.append(shape.text.strip())
    
    return "\n\n".join(text_parts)

def _get_speaker_notes(slide) -> str:
    """Extract speaker notes"""
    try:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if hasattr(notes_slide, 'notes_text_frame'):
                return notes_slide.notes_text_frame.text.strip()
    except:
        pass
    return ""

def get_deck_summary(slides: List[Dict]) -> Dict:
    """Get high-level deck statistics - FIXED FOR DIVISION BY ZERO"""
    if not slides:
        return {
            "total_slides": 0,
            "avg_content_length": 0,
            "slides_with_notes": 0,
            "slide_titles": [],
            "total_words": 0
        }
    
    total_content_length = sum(len(s['content']) for s in slides)
    
    return {
        "total_slides": len(slides),
        "avg_content_length": round(total_content_length / len(slides)),
        "slides_with_notes": sum(1 for s in slides if s['notes']),
        "slide_titles": [s['title'] for s in slides],
        "total_words": sum(len(s['content'].split()) for s in slides)
    }

def classify_deck_type(slides: List[Dict]) -> str:
    """Classify the type of pitch deck based on content"""
    if not slides:
        return "Unknown"
    
    all_text = " ".join([s['title'] + " " + s['content'] for s in slides]).lower()
    
    ai_ml_keywords = ['machine learning', 'deep learning', 'neural network', 'ai', 
                      'model', 'data', 'algorithm', 'prediction']
    
    keyword_count = sum(1 for keyword in ai_ml_keywords if keyword in all_text)
    
    if keyword_count >= 5:
        return "AI/ML Platform"
    elif keyword_count >= 3:
        return "Data Science Product"
    elif any(word in all_text for word in ['analytics', 'insights', 'dashboard']):
        return "Analytics Platform"
    else:
        return "Technology Company"